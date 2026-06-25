import io
import json
from datetime import datetime, timedelta

from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from pptx import Presentation
from sqlalchemy import desc, func
from sqlalchemy.orm import Session

from app.database import get_db
from app.deps import require_roles
from app.models import Analysis, Mention, ReportRun, User

router = APIRouter(prefix="/api/reports", tags=["reports"])


def _daily_snapshot(db: Session, hours: int = 24) -> dict:
    start_ts = datetime.utcnow() - timedelta(hours=min(max(hours, 1), 24 * 14))
    mention_count = db.query(Mention).filter(Mention.posted_at >= start_ts).count()
    harmful_count = (
        db.query(Analysis)
        .join(Mention, Mention.id == Analysis.mention_id)
        .filter(Mention.posted_at >= start_ts, Analysis.is_harmful == 1)
        .count()
    )
    avg_sentiment = (
        db.query(func.avg(Analysis.sentiment_score))
        .join(Mention, Mention.id == Analysis.mention_id)
        .filter(Mention.posted_at >= start_ts)
        .scalar()
    )
    top_topics = (
        db.query(Analysis.topic, func.count(Analysis.id).label("count"))
        .join(Mention, Mention.id == Analysis.mention_id)
        .filter(Mention.posted_at >= start_ts)
        .group_by(Analysis.topic)
        .order_by(desc("count"))
        .limit(5)
        .all()
    )
    return {
        "window_hours": hours,
        "mentions": mention_count,
        "harmful_alerts": harmful_count,
        "average_sentiment": round(float(avg_sentiment or 0), 3),
        "top_topics": [{"topic": topic, "count": count} for topic, count in top_topics],
        "generated_at": datetime.utcnow().isoformat(),
    }


def _record_report(db: Session, user: User, report_type: str, output_format: str, filters: dict) -> None:
    db.add(
        ReportRun(
            owner_email=user.email,
            report_type=report_type,
            format=output_format,
            filters_json=json.dumps(filters),
        )
    )
    db.commit()


@router.get("/daily")
def daily_report(
    hours: int = 24,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_roles("admin", "coordinator", "analyst")),
):
    data = _daily_snapshot(db, hours=hours)
    _record_report(db, current_user, "daily_intelligence", "json", {"hours": hours})
    return data


@router.get("/daily.pdf")
def daily_report_pdf(
    hours: int = 24,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_roles("admin", "coordinator", "analyst")),
):
    data = _daily_snapshot(db, hours=hours)
    _record_report(db, current_user, "daily_intelligence", "pdf", {"hours": hours})
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=A4)
    pdf.setTitle("Daily Intelligence Report")
    pdf.drawString(40, 800, "Narrative Intelligence - Daily Report")
    pdf.drawString(40, 780, f"Generated: {data['generated_at']}")
    pdf.drawString(40, 760, f"Mentions: {data['mentions']}")
    pdf.drawString(40, 740, f"Harmful Alerts: {data['harmful_alerts']}")
    pdf.drawString(40, 720, f"Average Sentiment: {data['average_sentiment']}")
    y = 700
    pdf.drawString(40, y, "Top Topics:")
    for row in data["top_topics"]:
        y -= 18
        pdf.drawString(55, y, f"- {row['topic']}: {row['count']}")
    pdf.showPage()
    pdf.save()
    buffer.seek(0)
    return StreamingResponse(buffer, media_type="application/pdf")


@router.get("/daily.pptx")
def daily_report_pptx(
    hours: int = 24,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_roles("admin", "coordinator", "analyst")),
):
    data = _daily_snapshot(db, hours=hours)
    _record_report(db, current_user, "daily_intelligence", "pptx", {"hours": hours})
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Daily Intelligence Report"
    slide.placeholders[1].text = f"Generated {data['generated_at']}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Metrics"
    text = slide2.placeholders[1].text_frame
    text.text = f"Mentions: {data['mentions']}"
    text.add_paragraph().text = f"Harmful Alerts: {data['harmful_alerts']}"
    text.add_paragraph().text = f"Average Sentiment: {data['average_sentiment']}"

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Top Topics"
    text2 = slide3.placeholders[1].text_frame
    text2.text = "Top Narrative Topics"
    for row in data["top_topics"]:
        text2.add_paragraph().text = f"{row['topic']}: {row['count']}"

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return StreamingResponse(
        out,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
